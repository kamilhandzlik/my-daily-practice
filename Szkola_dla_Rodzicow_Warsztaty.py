from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os

output_path = os.path.join(os.path.expanduser("~"), "Documents", "Szkola_dla_Rodzicow_Warsztaty.pptx")

# Create a presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Szkoła dla Rodziców"
subtitle.text = "Jak wychować dziecko i nie zwariować – świadome rodzicielstwo z sercem i strukturą"

# Slide 1: Wprowadzenie
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "I. Powitanie i integracja"
content = slide.placeholders[1]
content.text = (
    "- Przedstawienie prowadzącego i uczestników\n"
    "- Ćwiczenie: „Kim jestem jako rodzic?”\n"
    "- Wprowadzenie do programu warsztatów"
)

# Slide 2: Część 1
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "II. Dlaczego dzieci są „trudne”?"
content = slide.placeholders[1]
content.text = (
    "- Miniwykład: co kryje się za zachowaniami dziecka?\n"
    "- Rozwój emocjonalny i potrzeby dziecka\n"
    "- Ćwiczenie: „Zobacz świat oczami swojego dziecka”\n"
    "🧠 Hasło: „Nie złośliwość – tylko potrzeba!”"
)

# Slide 3: Część 2
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "III. Bezwarunkowa miłość vs. kontrola"
content = slide.placeholders[1]
content.text = (
    "- Znaczenie bezwarunkowej miłości\n"
    "- Granice jako wyraz miłości\n"
    "- Techniki: aktywne słuchanie, pauza, empatia\n"
    "- Ćwiczenie: „Trudna sytuacja – moja nowa reakcja”"
)

# Slide 4: Przerwa
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "☕️ Przerwa kawowa"
content = slide.placeholders[1]
content.text = "15 minut relaksu, rozmów i integracji"

# Slide 5: Część 3
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "IV. Tworzenie bezpiecznej przestrzeni"
content = slide.placeholders[1]
content.text = (
    "- Klimat emocjonalny w domu\n"
    "- Rytuały, język wdzięczności, rutyny\n"
    "- Praca z własnymi lękami\n"
    "- Ćwiczenie: „Nowy rytuał w naszym domu”"
)

# Slide 6: Część 4
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "V. Nauka przez doświadczenie"
content = slide.placeholders[1]
content.text = (
    "- Obserwacja jako metoda wychowawcza\n"
    "- Autorefleksja rodzicielska\n"
    "- Spokój i regulacja emocji"
)

# Slide 7: Podsumowanie
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "VI. Podsumowanie i zakończenie"
content = slide.placeholders[1]
content.text = (
    "- Feedback: „Co zabieram ze sobą?”\n"
    "- Rozdanie materiałów: karty ćwiczeń, lista książek\n"
    "- Zaproszenie do kolejnych warsztatów / grupy wsparcia"
)

# Slide 8: Materiały pomocnicze
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "📦 Materiały pomocnicze"
content = slide.placeholders[1]
content.text = (
    "- Mini-notes: „5 zasad świadomego rodzica”\n"
    "- Karta ćwiczeń: „Twój dom jako bezpieczna przystań”\n"
    "- Lista polecanych książek / podcastów\n"
    "- Drobny upominek (zakładka z cytatem)"
)

# Save presentation
prs.save(output_path)
output_path
